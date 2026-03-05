import os
import telebot
import requests
from flask import Flask, request
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from io import BytesIO
import random

app = Flask(__name__)

BOT_TOKEN = os.environ.get("BOT_TOKEN")
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")

bot = telebot.TeleBot(BOT_TOKEN, threaded=False)

genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel("gemini-pro")


# ---------------- IMAGE DOWNLOAD ----------------

def download_image(query):

    url = f"https://source.unsplash.com/800x600/?{query}"

    response = requests.get(url)

    return BytesIO(response.content)


# ---------------- PRESENTATION GENERATOR ----------------

def create_presentation(topic):

    prs = Presentation()

    slides_content = [
        f"{topic} haqida kirish",
        f"{topic} tushunchasi",
        f"{topic} tarixi",
        f"{topic} rivojlanishi",
        f"{topic} asosiy xususiyatlari",
        f"{topic} ishlash prinsipi",
        f"{topic} texnologiyalari",
        f"{topic} afzalliklari",
        f"{topic} kamchiliklari",
        f"{topic} qo'llanilish sohalari",
        f"{topic} dunyoda qo'llanilishi",
        f"{topic} O'zbekistonda qo'llanilishi",
        f"{topic} kelajak istiqbollari",
        f"{topic} qiziqarli faktlar",
        f"{topic} xulosa"
    ]

    for title in slides_content:

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = title

        img_stream = download_image(topic)

        image = Image.open(img_stream)

        img_path = f"{random.randint(1,9999)}.jpg"

        image.save(img_path)

        slide.shapes.add_picture(img_path, Inches(1), Inches(2), height=Inches(4))

        textbox = slide.shapes.add_textbox(Inches(5.5), Inches(2), Inches(4), Inches(4))
        text_frame = textbox.text_frame

        text_frame.text = f"{title} haqida muhim ma'lumotlar."

        for i in range(3):
            p = text_frame.add_paragraph()
            p.text = f"{topic} bo'yicha qo'shimcha tushuntirish {i+1}"
            p.level = 1

    filename = f"{topic}.pptx"

    prs.save(filename)

    return filename


# ---------------- START COMMAND ----------------

@bot.message_handler(commands=['start'])

def start(message):

    text = """
Salom!

Bu AI bot.

Siz:
• Savol berishingiz mumkin
• Yoki mavzu yozsangiz bot taqdimot yaratadi

Masalan:
Sun'iy intellekt
Python dasturlash
Kompyuter tarmoqlari
"""

    bot.send_message(message.chat.id, text)


# ---------------- CHAT HANDLER ----------------

@bot.message_handler(func=lambda message: True)

def chat(message):

    user_text = message.text

    try:

        bot.send_chat_action(message.chat.id, "typing")

        response = model.generate_content(user_text)

        if response.text:

            bot.reply_to(message, response.text)

        else:

            raise Exception("Gemini javob bermadi")

    except Exception as e:

        print("Gemini ishlamadi:", e)

        bot.send_message(message.chat.id, "Gemini ishlamadi. Taqdimot yaratilmoqda...")

        try:

            file = create_presentation(user_text)

            with open(file, "rb") as f:

                bot.send_document(message.chat.id, f)

        except Exception as e2:

            print("Presentation error:", e2)

            bot.send_message(message.chat.id, "Taqdimot yaratishda xatolik.")


# ---------------- WEBHOOK ----------------

@app.route("/", methods=["GET"])

def home():

    return "Bot ishlayapti"


@app.route("/", methods=["POST"])

def webhook():

    json_str = request.get_data().decode("UTF-8")

    update = telebot.types.Update.de_json(json_str)

    bot.process_new_updates([update])

    return "ok", 200


if __name__ == "__main__":

    port = int(os.environ.get("PORT", 10000))

    app.run(host="0.0.0.0", port=port)
